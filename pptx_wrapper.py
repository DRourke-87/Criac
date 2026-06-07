"""Generate professional .pptx presentations using python-pptx."""

from __future__ import annotations

import os
import tempfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ──────────────────────────────────────────────────────────
_BG       = RGBColor(0x11, 0x18, 0x27)  # Dark slate
_ACCENT   = RGBColor(0xF5, 0x9E, 0x0B)  # Amber gold
_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
_SILVER   = RGBColor(0xD1, 0xD5, 0xDB)  # Bullet text
_PANEL    = RGBColor(0x1F, 0x2A, 0x3C)  # Slightly lighter panel

# ── Canvas (16:9 widescreen) ────────────────────────────────────────────────
_W = Inches(13.333)
_H = Inches(7.5)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _text(slide, text: str, left, top, width, height,
          size: int, bold: bool = False, color: RGBColor = _WHITE,
          align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _title_slide(prs: Presentation, title: str) -> None:
    slide = _blank(prs)
    _bg(slide, _BG)
    _rect(slide, 0, 0, _W, Inches(0.12), _ACCENT)
    _rect(slide, 0, _H - Inches(0.12), _W, Inches(0.12), _ACCENT)
    _rect(slide, Inches(0.5), Inches(1.8), _W - Inches(1.0), Inches(3.9), _PANEL)
    _text(slide, title,
          Inches(1.0), Inches(2.2), _W - Inches(2.0), Inches(2.2),
          size=44, bold=True, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.0), Inches(4.45), Inches(3.33), Inches(0.05), _ACCENT)


def _content_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    slide = _blank(prs)
    _bg(slide, _BG)
    _rect(slide, 0, 0, Inches(0.1), _H, _ACCENT)
    _rect(slide, 0, 0, _W, Inches(1.1), _PANEL)
    _text(slide, heading,
          Inches(0.3), Inches(0.15), _W - Inches(0.6), Inches(0.85),
          size=30, bold=True)
    _rect(slide, Inches(0.3), Inches(1.15), _W - Inches(0.6), Inches(0.04), _ACCENT)

    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.35), _W - Inches(0.9), Inches(5.8)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"▸   {bullet}"
        run.font.size = Pt(20)
        run.font.color.rgb = _SILVER


def create_presentation(title: str, slides: list[dict]) -> str:
    """Build a .pptx file and return its /tmp path."""
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    _title_slide(prs, title)
    for s in slides[:7]:
        _content_slide(prs, s["heading"], s["bullets"])

    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:50].strip()
    path = os.path.join(tempfile.gettempdir(), f"{safe}.pptx")
    prs.save(path)
    return path
